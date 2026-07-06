"""
Servicios de importación: archivos → Django models.
Cada función recibe un objeto file-like y devuelve un resumen de lo procesado.
"""
import io
from datetime import datetime

import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation


# ─── PARSERS (archivo → datos crudos) ────────────────────────────────────────

def extraer_texto_pdf(archivo):
    """Extrae texto de PDF. Devuelve str."""
    texto = []
    with pdfplumber.open(archivo) as pdf:
        for pagina in pdf.pages:
            t = pagina.extract_text()
            if t:
                texto.append(t)
    return "\n".join(texto)


def extraer_tablas_pdf(archivo):
    """Extrae tablas de PDF como lista de DataFrames."""
    tablas = []
    with pdfplumber.open(archivo) as pdf:
        for pagina in pdf.pages:
            for tabla in pagina.extract_tables():
                if tabla:
                    df = pd.DataFrame(tabla[1:], columns=tabla[0])
                    tablas.append(df)
    return tablas


def extraer_texto_word(archivo):
    """Extrae texto de Word (.docx). Devuelve str."""
    doc = Document(archivo)
    parrafos = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(parrafos)


def extraer_tablas_word(archivo):
    """Extrae tablas de Word como lista de DataFrames."""
    doc = Document(archivo)
    tablas = []
    for tabla in doc.tables:
        filas = [[celda.text.strip() for celda in fila.cells] for fila in tabla.rows]
        if filas:
            df = pd.DataFrame(filas[1:], columns=filas[0])
            tablas.append(df)
    return tablas


def extraer_texto_pptx(archivo):
    """Extrae texto de PowerPoint (.pptx). Devuelve dict {slide_num: texto}."""
    prs = Presentation(archivo)
    slides = {}
    for i, slide in enumerate(prs.slides, start=1):
        textos = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                textos.append(shape.text.strip())
        slides[i] = "\n".join(textos)
    return slides


# ─── INGESTORES (datos → modelos Django) ─────────────────────────────────────

def excel_a_clientes(archivo):
    """
    Ingesta clientes desde Excel.

    Soporta el formato de importación original y el formato generado por la exportación.
    Devuelve dict con listas: creados, actualizados, errores.
    """
    from clientes.models import Cliente, PersonaNatural, PersonaJuridica
    import unicodedata

    df = pd.read_excel(archivo, dtype=str).fillna("")
    
    def clean_col(c):
        c = unicodedata.normalize('NFKD', str(c)).encode('ASCII', 'ignore').decode('utf-8')
        c = c.lower().strip()
        for char in [' ', '/', '-', '.']:
            c = c.replace(char, '_')
        while '__' in c:
            c = c.replace('__', '_')
        return c.strip('_')

    df.columns = [clean_col(c) for c in df.columns]

    creados, actualizados, errores = [], [], []

    for idx, fila in df.iterrows():
        tipo_raw = str(fila.get("tipo", "")).strip().lower()
        if "natural" in tipo_raw:
            tipo = "natural"
        elif any(term in tipo_raw for term in ["juridica", "jurídica", "empresa", "juridico", "jurídico"]):
            tipo = "juridica"
        else:
            tipo = tipo_raw


        email = str(fila.get("email", "")).strip()
        telefono = str(fila.get("telefono", "")).strip() or None
        
        activo_raw = str(fila.get("activo", "")).strip().lower()
        is_active = False if activo_raw in ("no", "false", "0", "f") else True
        
        num_fila = idx + 2  # 1-indexed + encabezado

        if not email:
            errores.append({"fila": num_fila, "error": "email vacío"})
            continue

        identificador = str(fila.get("identificador", "")).strip()
        nombre_rs = str(fila.get("nombre_razon_social", "")).strip()

        try:
            if tipo == "natural":
                run = str(fila.get("run", "")).strip() or identificador
                nombre = str(fila.get("nombre_completo", "")).strip() or nombre_rs
                if not run or not nombre:
                    errores.append({"fila": num_fila, "error": "run o nombre_completo vacío"})
                    continue

                try:
                    c = Cliente.objects.get(email_principal=email)
                    if not hasattr(c, 'personanatural'):
                        errores.append({"fila": num_fila, "error": "El email ya está registrado como Persona Jurídica"})
                        continue
                    pn = c.personanatural
                    pn.telefono_contacto = telefono
                    pn.is_active = is_active
                    pn.run = run
                    pn.nombre_completo = nombre
                    pn.save()
                    creado = False
                except Cliente.DoesNotExist:
                    PersonaNatural.objects.create(
                        email_principal=email,
                        telefono_contacto=telefono,
                        is_active=is_active,
                        run=run,
                        nombre_completo=nombre
                    )
                    creado = True

            elif tipo == "juridica":
                rut = str(fila.get("rut", "")).strip() or identificador
                razon = str(fila.get("razon_social", "")).strip() or nombre_rs
                giro = str(fila.get("giro", "")).strip() or "Sin giro"
                
                if not rut or not razon:
                    errores.append({"fila": num_fila, "error": "rut o razon_social vacío"})
                    continue

                try:
                    c = Cliente.objects.get(email_principal=email)
                    if not hasattr(c, 'personajuridica'):
                        errores.append({"fila": num_fila, "error": "El email ya está registrado como Persona Natural"})
                        continue
                    pj = c.personajuridica
                    pj.telefono_contacto = telefono
                    pj.is_active = is_active
                    pj.rut = rut
                    pj.razon_social = razon
                    pj.giro = giro
                    pj.save()
                    creado = False
                except Cliente.DoesNotExist:
                    PersonaJuridica.objects.create(
                        email_principal=email,
                        telefono_contacto=telefono,
                        is_active=is_active,
                        rut=rut,
                        razon_social=razon,
                        giro=giro
                    )
                    creado = True
            else:
                errores.append({"fila": num_fila, "error": f"tipo desconocido: '{tipo_raw}'"})
                continue

            if creado:
                creados.append(email)
            else:
                actualizados.append(email)

        except Exception as e:
            errores.append({"fila": num_fila, "error": str(e)})

    return {"creados": creados, "actualizados": actualizados, "errores": errores}


def excel_a_contratos(archivo):
    """
    Ingesta contratos desde Excel.

    Columnas esperadas:
        cliente_email | software_nombre | sla_nombre | tipo_contrato |
        monto | fecha_inicio | fecha_vencimiento | dias_gracia

    Devuelve dict con listas: creados, errores.
    """
    from clientes.models import Cliente
    from catalogo.models import Software
    from contratos.models import SLA, Contrato, TipoContrato

    df = pd.read_excel(archivo, dtype=str).fillna("")
    df.columns = [c.strip().lower().replace(" ", "_") for c in df.columns]

    creados, errores = [], []

    for idx, fila in df.iterrows():
        num_fila = idx + 2
        try:
            cliente = Cliente.objects.get(email_principal=fila["cliente_email"].strip())
            software = Software.objects.get(nombre=fila["software_nombre"].strip())
            sla = SLA.objects.get(nombre=fila["sla_nombre"].strip())

            tipo_raw = fila.get("tipo_contrato", "").strip().upper()
            tipos_validos = {t.value: t for t in TipoContrato}
            if tipo_raw not in tipos_validos:
                errores.append({"fila": num_fila, "error": f"tipo_contrato inválido: '{tipo_raw}'"})
                continue

            fecha_inicio = pd.to_datetime(fila["fecha_inicio"]).date()
            fecha_venc_raw = fila.get("fecha_vencimiento", "").strip()
            fecha_venc = pd.to_datetime(fecha_venc_raw).date() if fecha_venc_raw else None
            monto = float(fila.get("monto", 0) or 0)
            dias_gracia = int(fila.get("dias_gracia", 0) or 0)

            contrato = Contrato.objects.create(
                cliente=cliente,
                software=software,
                sla=sla,
                tipo_contrato=tipo_raw,
                monto=monto,
                fecha_inicio=fecha_inicio,
                fecha_vencimiento=fecha_venc,
                dias_gracia_autorizados=dias_gracia,
            )
            creados.append(contrato.id)

        except Cliente.DoesNotExist:
            errores.append({"fila": num_fila, "error": f"cliente no encontrado: {fila.get('cliente_email')}"})
        except Software.DoesNotExist:
            errores.append({"fila": num_fila, "error": f"software no encontrado: {fila.get('software_nombre')}"})
        except SLA.DoesNotExist:
            errores.append({"fila": num_fila, "error": f"SLA no encontrado: {fila.get('sla_nombre')}"})
        except Exception as e:
            errores.append({"fila": num_fila, "error": str(e)})

    return {"creados": creados, "errores": errores}
